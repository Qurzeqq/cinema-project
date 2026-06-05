from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "presentation_assets"
OUT = ROOT / "Vorobyev_defense_presentation.pptx"

W = Inches(13.333)
H = Inches(7.5)

BG = RGBColor(255, 249, 251)
TEXT = RGBColor(75, 47, 58)
MUTED = RGBColor(126, 92, 105)
ACCENT = RGBColor(236, 153, 177)
ACCENT_DARK = RGBColor(139, 70, 89)
MINT = RGBColor(190, 231, 205)
BLUE = RGBColor(63, 125, 164)
ORANGE = RGBColor(241, 178, 103)
LILAC = RGBColor(208, 190, 231)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(239, 207, 217)


def inch(value):
    return Inches(value)


def set_font(run, size=18, color=TEXT, bold=False):
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def fill_background(slide, color=BG):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def add_footer(slide, idx, total, dark=False):
    line_color = RGBColor(70, 58, 64) if dark else LINE
    text_color = RGBColor(214, 207, 211) if dark else MUTED
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(0.55), inch(7.05), inch(12.2), inch(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = line_color
    line.line.fill.background()
    box = slide.shapes.add_textbox(inch(0.58), inch(7.12), inch(12.0), inch(0.25))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"Воробьев К. Д. | Защита дипломной работы | {idx}/{total}"
    set_font(run, 8.5, text_color)


def add_title(slide, title, kicker=None):
    if kicker:
        tag = slide.shapes.add_textbox(inch(0.65), inch(0.28), inch(2.6), inch(0.28))
        tf = tag.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = kicker.upper()
        set_font(run, 8.5, ACCENT_DARK, bold=True)
    box = slide.shapes.add_textbox(inch(0.62), inch(0.52), inch(12.1), inch(0.72))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_font(run, 27, TEXT, bold=True)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(0.62), inch(1.25), inch(1.05), inch(0.07))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()


def add_text(slide, text, x, y, w, h, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = inch(0.02)
    tf.margin_right = inch(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size, color, bold)
    return box


def add_bullets(slide, bullets, x, y, w, h, size=17, color=TEXT, bullet_color=None):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = inch(0.02)
    tf.margin_right = inch(0.02)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(7)
        p.line_spacing = 1.08
        run_b = p.add_run()
        run_b.text = "• "
        set_font(run_b, size, bullet_color or ACCENT_DARK, bold=True)
        run = p.add_run()
        run.text = item
        set_font(run, size, color)
    return box


def add_card(slide, x, y, w, h, title, body=None, fill=WHITE, border=LINE, title_size=15, body_size=12):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(w), inch(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(0.8)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = inch(0.16)
    tf.margin_right = inch(0.16)
    tf.margin_top = inch(0.12)
    tf.margin_bottom = inch(0.08)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.space_after = Pt(4)
    r = p.add_run()
    r.text = title
    set_font(r, title_size, ACCENT_DARK, bold=True)
    if body:
        p2 = tf.add_paragraph()
        p2.line_spacing = 1.04
        r2 = p2.add_run()
        r2.text = body
        set_font(r2, body_size, TEXT)
    return shape


def add_chip(slide, x, y, text, color=ACCENT, w=None):
    w = w or max(1.1, 0.16 * len(text) + 0.36)
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(w), inch(0.36))
    chip.fill.solid()
    chip.fill.fore_color.rgb = color
    chip.line.fill.background()
    tf = chip.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    set_font(r, 10.5, TEXT, bold=True)
    return chip


def add_image(slide, image_name, x, y, w, h, border=True):
    path = ASSETS / image_name
    if not path.exists():
        raise FileNotFoundError(path)
    if border:
        frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(w), inch(h))
        frame.fill.solid()
        frame.fill.fore_color.rgb = WHITE
        frame.line.color.rgb = LINE
        frame.line.width = Pt(0.8)
    with Image.open(path) as img:
        iw, ih = img.size
    scale = min(inch(w) / iw, inch(h) / ih)
    pw, ph = int(iw * scale), int(ih * scale)
    px = inch(x) + int((inch(w) - pw) / 2)
    py = inch(y) + int((inch(h) - ph) / 2)
    pic = slide.shapes.add_picture(str(path), px, py, width=pw, height=ph)
    return pic


def add_flow(slide, labels, x, y, w, color=ACCENT):
    gap = 0.1
    cell_w = (w - gap * (len(labels) - 1)) / len(labels)
    for i, label in enumerate(labels):
        add_card(slide, x + i * (cell_w + gap), y, cell_w, 0.55, label, fill=RGBColor(255, 240, 245), title_size=10.5)
        if i < len(labels) - 1:
            add_text(slide, "→", x + (i + 1) * cell_w + i * gap - 0.02, y + 0.12, 0.18, 0.2, 14, color, True, PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]
    slides = []

    # 1
    slide = prs.slides.add_slide(blank)
    fill_background(slide, RGBColor(18, 18, 20))
    add_image(slide, "image2.png", 1.05, 0.35, 11.2, 2.1, border=False)
    add_text(
        slide,
        "Создание системы онлайн-продажи билетов\nи управления бронированием для сети кинотеатров",
        0.85,
        2.85,
        11.8,
        1.3,
        31,
        WHITE,
        True,
        PP_ALIGN.CENTER,
    )
    add_text(slide, "Дипломная работа", 0.85, 4.28, 11.8, 0.35, 18, ACCENT, True, PP_ALIGN.CENTER)
    add_card(
        slide,
        1.4,
        4.9,
        4.2,
        0.92,
        "Студент",
        "Воробьев Константин Дмитриевич\nгруппа 4 ИСП9-33",
        fill=RGBColor(255, 244, 247),
        border=RGBColor(92, 69, 78),
    )
    add_card(
        slide,
        7.7,
        4.9,
        4.2,
        0.92,
        "Специальность",
        "09.02.07 Информационные системы и программирование",
        fill=RGBColor(255, 244, 247),
        border=RGBColor(92, 69, 78),
    )
    add_text(slide, "Москва, 2026", 0.85, 6.25, 11.8, 0.35, 15, RGBColor(215, 210, 213), False, PP_ALIGN.CENTER)
    slides.append(slide)

    # 2
    slide = prs.slides.add_slide(blank)
    fill_background(slide)
    add_title(slide, "Актуальность и проблема", "01")
    add_text(
        slide,
        "Кинотеатру нужна единая цифровая среда, где расписание, места, бронирования и продажи обрабатываются согласованно.",
        0.75,
        1.55,
        5.6,
        0.75,
        18,
        TEXT,
        True,
    )
    add_bullets(
        slide,
        [
            "рост количества онлайн-заказов повышает нагрузку на персонал;",
            "разрозненные инструменты учета приводят к ошибкам и потере актуальности данных;",
            "для клиента критичны быстрый выбор сеанса, места и оплата с любого устройства;",
            "для сети кинотеатров важны централизованное управление и контроль загрузки залов.",
        ],
        0.78,
        2.45,
        5.55,
        3.35,
        16.5,
    )
    add_card(slide, 6.8, 1.55, 2.75, 1.2, "Ошибки", "двойные брони и ручная обработка заказов", fill=RGBColor(255, 239, 242))
    add_card(slide, 9.85, 1.55, 2.75, 1.2, "Разрозненность", "нет единой базы для клиентов, фильмов и сеансов", fill=RGBColor(240, 248, 255))
    add_card(slide, 6.8, 3.1, 2.75, 1.2, "Задержки", "неактуальная информация о свободных местах", fill=RGBColor(245, 250, 239))
    add_card(slide, 9.85, 3.1, 2.75, 1.2, "Оплата", "нужен удобный переход от брони к покупке", fill=RGBColor(255, 247, 235))
    add_image(slide, "image3.png", 6.9, 4.75, 5.55, 1.28, border=True)
    slides.append(slide)

    # 3
    slide = prs.slides.add_slide(blank)
    fill_background(slide)
    add_title(slide, "Цель и задачи работы", "02")
    add_card(
        slide,
        0.75,
        1.52,
        12.0,
        1.05,
        "Цель дипломной работы",
        "разработка информационной системы, обеспечивающей автоматизацию онлайн-продажи билетов и управления бронированием мест для сети кинотеатров.",
        fill=RGBColor(255, 239, 245),
        title_size=16,
        body_size=13,
    )
    add_card(slide, 0.75, 2.9, 3.65, 1.1, "Объект", "процессы автоматизации онлайн-продаж и бронирования билетов в сфере киноиндустрии.", fill=WHITE)
    add_card(slide, 4.82, 2.9, 3.65, 1.1, "Предмет", "ПО для продажи билетов, бронирования мест, управления расписанием и аналитики.", fill=WHITE)
    add_card(slide, 8.9, 2.9, 3.65, 1.1, "Результат", "работоспособное веб-приложение с пользовательской и административной частью.", fill=WHITE)
    add_bullets(
        slide,
        [
            "проанализировать предметную область и существующие решения;",
            "сформировать функциональные и нефункциональные требования;",
            "спроектировать архитектуру, базу данных и пользовательские сценарии;",
            "реализовать модули бронирования, продажи билетов и администрирования;",
            "провести тестирование и оценить корректность работы системы.",
        ],
        0.95,
        4.55,
        11.4,
        1.8,
        16,
    )
    slides.append(slide)

    # 4
    slide = prs.slides.add_slide(blank)
    fill_background(slide)
    add_title(slide, "Бизнес-процессы и требования", "03")
    add_image(slide, "image6.png", 0.68, 1.55, 7.1, 3.48, border=True)
    add_image(slide, "image7.png", 7.98, 1.55, 4.85, 2.82, border=True)
    add_card(slide, 8.0, 4.68, 1.45, 0.72, "A1", "пользователи", fill=RGBColor(237, 246, 255), title_size=13, body_size=10)
    add_card(slide, 9.62, 4.68, 1.45, 0.72, "A2", "бронирование", fill=RGBColor(239, 252, 241), title_size=13, body_size=10)
    add_card(slide, 11.24, 4.68, 1.45, 0.72, "A3", "администрирование", fill=RGBColor(255, 246, 232), title_size=13, body_size=9)
    add_bullets(
        slide,
        [
            "пользовательский контур: афиша, выбор сеанса, схема зала, бронирование, оплата;",
            "административный контур: управление фильмами, сеансами и бронированиями;",
            "информационное обеспечение: хранение данных, проверка доступности мест, QR-код и отправка билета.",
        ],
        0.85,
        5.45,
        11.7,
        1.05,
        13.5,
    )
    slides.append(slide)

    # 5
    slide = prs.slides.add_slide(blank)
    fill_background(slide)
    add_title(slide, "Архитектура и стек технологий", "04")
    add_image(slide, "image42.png", 0.72, 1.48, 7.3, 4.12, border=True)
    add_bullets(
        slide,
        [
            "клиентская часть работает в браузере пользователя или администратора;",
            "серверная часть развёрнута в локальной среде XAMPP: Apache + PHP;",
            "данные хранятся в MySQL, взаимодействие с БД выполнено через PDO;",
            "внешние сервисы в демонстрационной версии: QR-страница оплаты и отправка email-билета.",
        ],
        8.33,
        1.65,
        4.1,
        2.45,
        15,
    )
    for i, (label, color) in enumerate(
        [
            ("PHP", ACCENT),
            ("MySQL", MINT),
            ("Apache", ORANGE),
            ("HTML/CSS", LILAC),
            ("JavaScript", RGBColor(247, 222, 130)),
            ("PDO", RGBColor(208, 226, 244)),
        ]
    ):
        add_chip(slide, 8.4 + (i % 2) * 1.65, 4.4 + (i // 2) * 0.5, label, color, w=1.35)
    slides.append(slide)

    # 6
    slide = prs.slides.add_slide(blank)
    fill_background(slide)
    add_title(slide, "Структура базы данных", "05")
    add_image(slide, "image11.png", 5.15, 1.48, 7.42, 4.35, border=True)
    add_card(slide, 0.78, 1.55, 3.95, 0.8, "clients", "пользователи, email, пароль, роль", fill=RGBColor(239, 247, 255))
    add_card(slide, 0.78, 2.53, 3.95, 0.8, "films", "название, описание, длительность, цена", fill=RGBColor(255, 244, 247))
    add_card(slide, 0.78, 3.51, 3.95, 0.8, "sessions", "фильм, дата, время, зал", fill=RGBColor(247, 252, 239))
    add_card(slide, 0.78, 4.49, 3.95, 0.8, "booking", "клиент, сеанс, ряд, место, дата брони", fill=RGBColor(255, 248, 235))
    add_bullets(
        slide,
        [
            "связи «один ко многим»: фильм → сеансы, сеанс → бронирования, клиент → бронирования;",
            "структура позволяет получать полные данные о билете: пользователь, фильм, время и места;",
            "проверка занятости мест выполняется перед созданием записи бронирования.",
        ],
        0.95,
        5.78,
        11.5,
        0.9,
        12.8,
    )
    slides.append(slide)

    # 7
    slide = prs.slides.add_slide(blank)
    fill_background(slide)
    add_title(slide, "Пользовательский сценарий", "06")
    add_image(slide, "image13.png", 0.7, 1.48, 5.9, 4.05, border=True)
    add_image(slide, "image21.png", 6.85, 1.48, 5.78, 4.05, border=True)
    add_flow(
        slide,
        ["поиск фильма", "выбор сеанса", "схема зала", "подтверждение", "переход к оплате"],
        0.82,
        5.85,
        11.55,
    )
    slides.append(slide)

    # 8
    slide = prs.slides.add_slide(blank)
    fill_background(slide)
    add_title(slide, "Логика бронирования", "07")
    add_image(slide, "image40.png", 0.7, 1.45, 7.65, 3.72, border=True)
    add_bullets(
        slide,
        [
            "пользователь выбирает сеанс и свободные места на схеме зала;",
            "сервер проверяет авторизацию и доступность выбранных мест;",
            "при успешной проверке создаются записи бронирования и выполняется переход к оплате;",
            "занятые места блокируются для повторного выбора, бронь можно отменить из личного кабинета.",
        ],
        8.68,
        1.65,
        3.75,
        2.75,
        15,
    )
    add_card(
        slide,
        8.75,
        4.75,
        3.55,
        0.95,
        "Ключевой эффект",
        "снижение риска двойного бронирования и сохранение целостности данных о местах.",
        fill=RGBColor(239, 252, 241),
        title_size=14,
        body_size=11,
    )
    slides.append(slide)

    # 9
    slide = prs.slides.add_slide(blank)
    fill_background(slide)
    add_title(slide, "Оплата и электронный билет", "08")
    add_image(slide, "image26.png", 0.72, 1.5, 6.45, 4.15, border=True)
    add_bullets(
        slide,
        [
            "после подтверждения выбранных мест пользователь попадает на страницу оплаты;",
            "на странице отображаются фильм, дата, время, зал, места, количество билетов и итоговая сумма;",
            "QR-код открывает демонстрационную страницу оплаты;",
            "после подтверждения система формирует электронный билет и отправляет его на почту пользователя.",
        ],
        7.6,
        1.75,
        4.85,
        2.6,
        15,
    )
    add_card(
        slide,
        7.72,
        4.7,
        4.55,
        0.9,
        "Ограничение текущей версии",
        "модуль оплаты демонстрационный; реальная платёжная интеграция запланирована как направление развития.",
        fill=RGBColor(255, 247, 235),
        title_size=14,
        body_size=11,
    )
    slides.append(slide)

    # 10
    slide = prs.slides.add_slide(blank)
    fill_background(slide)
    add_title(slide, "Административная панель", "09")
    add_image(slide, "image29.png", 0.68, 1.48, 7.7, 4.78, border=True)
    add_card(slide, 8.65, 1.62, 3.75, 0.75, "Статистика", "фильмы, сеансы, бронирования, пользователи", fill=RGBColor(255, 244, 247))
    add_card(slide, 8.65, 2.58, 3.75, 0.75, "Фильмы", "добавление, редактирование и удаление", fill=RGBColor(239, 247, 255))
    add_card(slide, 8.65, 3.54, 3.75, 0.75, "Сеансы", "создание расписания и изменение данных", fill=RGBColor(247, 252, 239))
    add_card(slide, 8.65, 4.5, 3.75, 0.75, "Бронирования", "просмотр занятых мест и списка броней", fill=RGBColor(255, 248, 235))
    add_bullets(slide, ["доступ к панели ограничен ролью администратора."], 8.78, 5.62, 3.55, 0.45, 13.5)
    slides.append(slide)

    # 11
    slide = prs.slides.add_slide(blank)
    fill_background(slide)
    add_title(slide, "Тестирование приложения", "10")
    add_image(slide, "image36.png", 0.72, 1.5, 6.7, 3.14, border=True)
    add_bullets(
        slide,
        [
            "синтаксическая проверка PHP-файлов проекта;",
            "проверка подключения к БД и наличия таблиц clients, films, sessions, booking;",
            "проверка доступности основных страниц и защиты закрытых разделов;",
            "сквозной сценарий пользователя: регистрация, бронирование, переход к оплате;",
            "административный сценарий: вход, создание фильма и сеанса.",
        ],
        7.82,
        1.58,
        4.6,
        2.7,
        14.5,
    )
    add_card(slide, 0.82, 5.0, 2.65, 0.78, "19 PHP-файлов", "синтаксическая проверка без ошибок", fill=RGBColor(239, 252, 241), title_size=13, body_size=10)
    add_card(slide, 3.72, 5.0, 2.65, 0.78, "HTTP 200", "главная, вход, регистрация, бронирование", fill=RGBColor(239, 247, 255), title_size=13, body_size=10)
    add_card(slide, 6.62, 5.0, 2.65, 0.78, "302 redirect", "защита закрытых страниц", fill=RGBColor(255, 248, 235), title_size=13, body_size=10)
    add_card(slide, 9.52, 5.0, 2.65, 0.78, "ALL TESTS PASSED", "функциональная проверка успешна", fill=RGBColor(255, 244, 247), title_size=13, body_size=10)
    slides.append(slide)

    # 12
    slide = prs.slides.add_slide(blank)
    fill_background(slide)
    add_title(slide, "Итоги и развитие", "11")
    add_card(
        slide,
        0.78,
        1.5,
        5.75,
        1.0,
        "Результат работы",
        "разработана информационная система, которая поддерживает полный пользовательский сценарий бронирования билета и административное управление данными.",
        fill=RGBColor(255, 239, 245),
        title_size=16,
        body_size=12,
    )
    add_bullets(
        slide,
        [
            "пользователь может найти фильм, выбрать сеанс и места, оформить бронь и перейти к оплате;",
            "администратор управляет фильмами, сеансами и просматривает список бронирований;",
            "БД обеспечивает хранение связанных данных о клиентах, фильмах, сеансах и билетах;",
            "тестирование подтвердило корректность основных функций системы.",
        ],
        0.92,
        2.85,
        5.4,
        2.2,
        15,
    )
    add_card(slide, 7.0, 1.5, 5.55, 0.72, "Направления развития", fill=RGBColor(239, 247, 255), title_size=16)
    add_card(slide, 7.0, 2.48, 2.55, 0.88, "Платежи", "интеграция с реальным платёжным сервисом", fill=WHITE, title_size=13, body_size=10)
    add_card(slide, 9.95, 2.48, 2.55, 0.88, "Статусы", "автоматическая проверка оплаты и управление жизненным циклом", fill=WHITE, title_size=13, body_size=10)
    add_card(slide, 7.0, 3.68, 2.55, 0.88, "Отчётность", "аналитика загрузки залов и продаж", fill=WHITE, title_size=13, body_size=10)
    add_card(slide, 9.95, 3.68, 2.55, 0.88, "Масштабирование", "поддержка нескольких кинотеатров и залов", fill=WHITE, title_size=13, body_size=10)
    add_text(slide, "Цель дипломной работы достигнута.", 0.82, 5.75, 11.8, 0.55, 25, ACCENT_DARK, True, PP_ALIGN.CENTER)
    add_text(slide, "Спасибо за внимание", 0.82, 6.25, 11.8, 0.35, 18, TEXT, False, PP_ALIGN.CENTER)
    slides.append(slide)

    total = len(slides)
    for idx, slide in enumerate(slides, start=1):
        add_footer(slide, idx, total, dark=(idx == 1))

    prs.save(OUT)
    print(OUT)
    print(f"slides={len(slides)}")


if __name__ == "__main__":
    build()
